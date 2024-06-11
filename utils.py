import base64
import mimetypes
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.placeholder import PlaceholderPicture

class Feedback:
    def __init__(self, slide_number):
        self.slide_number = slide_number
        self.feedback = []
        self.adjusted_images = []
    def add_feedback(self, fb):
        self.feedback.append(fb)
    def add_image(self, img):
        self.adjusted_images.append(img)
    def format_json(self):
        return {
            "slide": self.slide_number,
            "feedback": "\n".join(list(set(self.feedback))),
            "images": self.adjusted_images
        }


def is_image(shape):
    return shape.shape_type in [MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.PICTURE]
def get_image(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        if isinstance(shape, PlaceholderPicture):
            return shape.image
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            return get_image(s)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return shape.image

def write_image(image, filename):
    with open(filename, "wb") as f:
        f.write(image.blob)
    return filename

def get_image_uri(filename):
    mime, _ = mimetypes.guess_type(filename)
    with open(filename, "rb") as file:
        data = file.read()
        d = base64.encodebytes(data).splitlines()
        data64 = u''
        for i in d:
            data64 = data64 + i.decode()
        return u'data:%s;base64,%s' % (mime, data64)

