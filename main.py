from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from utils import is_image, get_image, write_image

p = Presentation("QRL.pptx")
i = 0

for slide in p.slides:
    print(slide.follow_master_background)
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            # print(text_frame.margin_top.inches)
            # print(text_frame.margin_bottom.inches)
            # print(text_frame.margin_left.inches)
            # print(text_frame.margin_right.inches)
            # for paragraph in text_frame.paragraphs:
            #     print(paragraph.font)
            #     print(paragraph.text)
            #     print(paragraph.level)
        elif is_image(shape):
            i += 1
            write_image(get_image(shape), f"{i}.jpg")
