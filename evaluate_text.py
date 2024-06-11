from pptx.enum.text import PP_ALIGN

def evaluate_text(slide, fb):
    check_fonts(slide, fb)
    check_alignment(slide, fb)
    check_visibility(slide, fb)

def check_fonts(slide, fb):
    title_font = slide.shapes.title.text_frame.paragraphs[0].font.name
    if title_font not in ["Mecherle Sans Semibold", "Mecherle Sans Slab Regular"]:
        fb.add_feedback(f"Consider changing the font of the title '{slide.shapes.title}' from {title_font} to Mecherle Sans Semibold or Mecherle Sans Slab Regular")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if paragraph.level == 0:
                    if paragraph.font.name not in ["Mecherle Sans Semibold", "Mecherle Sans Slab Regular"]:
                        fb.add_feedback(f"Consider changing the font of '{paragraph.text}' from {paragraph.font.name} to Mecherle Sans Semibold or Mecherle Sans Slab Regular")
                else:
                    if paragraph.font.name not in ["Mecherle Sans Medium", "Mecherle Sans Regular", "Mecherle Sans Semibold"]:
                        fb.add_feedback(f"Consider changing the font of '{paragraph.text}' from {paragraph.font.name} to Mecherle Sans Medium, Regular, or Semibold")


def check_alignment(slide, fb):
    title_align = slide.shapes.title.text_frame.paragraphs[0].alignment
    if title_align not in [PP_ALIGN.CENTER, PP_ALIGN.LEFT]:
        fb.add_feedback("Titles should be left or center aligned")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if paragraph.level == 0:
                    if paragraph.alignment == PP_ALIGN.CENTER:
                        fb.add_feedback("Ensure that center alignment is being used in central composition")
                    elif paragraph.alignment != PP_ALIGN.LEFT:
                        fb.add_feedback("Text should typically be left-aligned")


def check_visibility(slide, fb):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if paragraph.font.size < 12:
                    fb.add_feedback(f"Consider increasing the font size from {paragraph.font.size} to at least 12")
                # if paragraph.font.color.rgb:
                    # do color stuff here